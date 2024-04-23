import unittest
from unittest.mock import Mock
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
from to_markdown import PPTXtoMarkdown  # Ensure this matches the actual script name

class TestPPTXtoMarkdown(unittest.TestCase):

    def setUp(self):
        # Mocking the paths
        self.pptx_path = "dummy_path.pptx"
        self.md_path = "dummy_output.md"
        
        # Create an instance of the converter
        self.converter = PPTXtoMarkdown(self.pptx_path, self.md_path)

    def test_process_text_frame_with_body_text(self):
        # Mocking the shape object and its properties
        shape = Mock()
        shape.text_frame = Mock()
        shape.text_frame.text = "Test Body Text"
        
        shape.placeholder_format = Mock()
        shape.placeholder_format.type = Mock()
        shape.placeholder_format.type.name = "BODY"  # Proper enum simulation
        
        # Mocking the markdown file handle
        md_file = Mock()
        
        # Run the method
        self.converter.process_text_frame(md_file, shape)
        
        # Check if the markdown file's write method was called correctly
        md_file.write.assert_called_with("Test Body Text\n\n")
    
    def test_process_text_frame_with_code_example(self):
        # Mocking the shape object and its properties
        shape = Mock()
        shape.text_frame = Mock()
        shape.text_frame.text = "print('Hello, world')"
        
        shape.fill = Mock()
        shape.fill.fore_color = Mock()
        shape.fill.fore_color.type = MSO_COLOR_TYPE.RGB
        shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        
        shape.placeholder_format = Mock()
        shape.placeholder_format.type = Mock()
        shape.placeholder_format.type.name = "TITLE"  # Ensure this does not trigger 'body' text processing
        
        # Mocking the markdown file handle
        md_file = Mock()
        
        # Run the method
        self.converter.process_text_frame(md_file, shape)
        
        # Check if the markdown file's write method was called correctly for code example
        md_file.write.assert_called_with("```python\nprint('Hello, world')\n```\n\n")

    def test_is_code_example_with_black_solid_fill(self):
        # Mocking the shape object and its fill properties
        shape = Mock()
        shape.fill = Mock()
        shape.fill.solid = Mock(return_value=True)
        shape.fill.fore_color = Mock()
        shape.fill.fore_color.type = MSO_COLOR_TYPE.RGB
        shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        
        # Run the method
        result = self.converter.is_code_example(shape)
        
        # Check if the method returns True for black solid fill
        self.assertTrue(result)

    def test_is_code_example_with_none_fill(self):
        # Mocking the shape object with _NoneFill type which should return False
        shape = Mock()
        shape.fill = Mock()
        shape.fill.solid = Mock(return_value=False)
        
        # Run the method
        result = self.converter.is_code_example(shape)
        
        # Check if the method returns False for non-solid fill
        self.assertFalse(result)


if __name__ == '__main__':
    unittest.main()