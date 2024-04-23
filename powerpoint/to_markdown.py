import pptx
import argparse
import os


class PPTXtoMarkdown:
    def __init__(self, pptx_path, md_path):
        self.pptx_path = pptx_path
        self.md_path = md_path

    def extract_content(self):
        try:
            presentation = pptx.Presentation(self.pptx_path)
        except FileNotFoundError:
            print("Error: PowerPoint file not found.")
            return
        except IOError:
            print("Error: PowerPoint file could not be opened.")
            return

        with open(self.md_path, 'w', encoding='utf-8') as md_file:
            for slide_number, slide in enumerate(presentation.slides, start=1):
                self.process_slide(md_file, slide_number, slide)

    def process_slide(self, md_file, slide_number, slide):
        md_file.write(f"## Slide {slide_number}\n\n")

        if slide.shapes.title:
            md_file.write(f"### {slide.shapes.title.text}\n\n")

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                self.process_text_frame(md_file, shape)

        if slide.has_notes_slide:
            self.write_speaker_notes(md_file, slide)

        md_file.write("---\n\n")

    def process_text_frame(self, md_file, shape):
        if hasattr(shape, 'text_frame') and shape.text_frame and hasattr(shape.text_frame, 'text'):
            text = shape.text_frame.text.strip()
            if text:
                # Ensure shape is a placeholder before checking type
                if (hasattr(shape, 'is_placeholder') and shape.is_placeholder and
                    shape.placeholder_format and shape.placeholder_format.type):
                    if 'body' in shape.placeholder_format.type.name.lower():
                        md_file.write(f"{text}\n\n")
                
                if self.is_code_example(shape):
                    md_file.write(f"```python\n{text}\n```\n\n")

    def is_code_example(self, shape):
        if hasattr(shape, 'fill') and shape.fill:
            if shape.fill.solid():
                return (shape.fill.fore_color and
                        shape.fill.fore_color.type == pptx.enum.dml.MSO_COLOR_TYPE.RGB and
                        shape.fill.fore_color.rgb == pptx.dml.color.RGBColor(0, 0, 0))
        return False

    def write_speaker_notes(self, md_file, slide):
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text.strip() if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text else ''
        if notes_text:
            md_file.write(f"**Speaker Notes:**\n\n{notes_text}\n\n")


def parse_arguments():
    parser = argparse.ArgumentParser(description="Extracts content from PPTX and outputs to Markdown format.")
    parser.add_argument("pptx_path", help="Path to the PowerPoint file.")
    parser.add_argument("output_md_path", help="Path for the output Markdown file.")
    args = parser.parse_args()
    return args.pptx_path, args.output_md_path


if __name__ == "__main__":
    pptx_path, md_path = parse_arguments()
    converter = PPTXtoMarkdown(pptx_path, md_path)
    converter.extract_content()
    print("Content extraction completed. Check the Markdown file.")
